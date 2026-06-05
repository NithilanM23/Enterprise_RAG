"""
services/metadata_service.py
-----------------------------
Document metadata extraction and storage for all formats.

Stores structural metadata about each document at ingestion time so that
metadata questions ("how many pages?", "what columns?", "what sheets?")
can be answered instantly without going through the RAG pipeline.

Metadata stored per document:
  PDF   → page_count, file_size_kb, word_count (approx)
  DOCX  → page_count (approx), word_count, table_count
  PPTX  → slide_count, has_images, file_size_kb
  TXT   → line_count, word_count, file_size_kb
  XLSX  → sheet_names, row_counts per sheet, column_names per sheet

Schema:
  document_metadata
    id, document_id, key, value
    One row per metadata key-value pair per document.

Metadata questions are detected by keywords:
  "how many", "pages", "rows", "columns", "sheets", "count",
  "size", "structure", "fields", "headers", "what columns"
"""

import logging
import os
import sys
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Metadata question detection keywords
# ---------------------------------------------------------------------------

METADATA_KEYWORDS = {
    "pages", "page", "slides", "slide", "rows", "row", "columns", "column",
    "sheets", "sheet", "fields", "field", "headers", "header", "structure",
    "size", "count", "many", "total", "metadata", "info", "about",
}


def is_metadata_question(query: str) -> bool:
    """
    Detect if a query is asking about document structure rather than content.

    Returns True if the query contains metadata-related keywords.
    Examples:
        "how many pages does this document have?" → True
        "what columns are in the sales sheet?" → True
        "where is the company located?" → False
    """
    tokens = set(query.lower().split())
    return bool(tokens & METADATA_KEYWORDS)


# ---------------------------------------------------------------------------
# Table creation
# ---------------------------------------------------------------------------

def ensure_metadata_table() -> None:
    """Create document_metadata table if it does not exist. Idempotent."""
    from services.database_service import _get_raw_connection

    with _get_raw_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE TABLE IF NOT EXISTS document_metadata (
                    id          SERIAL PRIMARY KEY,
                    document_id INTEGER NOT NULL
                                    REFERENCES documents(id) ON DELETE CASCADE,
                    key         TEXT NOT NULL,
                    value       TEXT NOT NULL,
                    CONSTRAINT document_metadata_unique UNIQUE (document_id, key)
                );
            """)

    logger.debug("document_metadata table ensured.")


# ---------------------------------------------------------------------------
# Metadata extraction per format
# ---------------------------------------------------------------------------

def extract_metadata(filepath: str, document_id: int) -> dict:
    """
    Extract structural metadata from a document and store it in the DB.

    Args:
        filepath    : Path to the document file.
        document_id : The document's DB row ID.

    Returns:
        Dict of key → value metadata pairs that were stored.
    """
    path = Path(filepath)
    ext  = path.suffix.lower()

    extractors = {
        ".pdf":  _extract_pdf_metadata,
        ".docx": _extract_docx_metadata,
        ".pptx": _extract_pptx_metadata,
        ".ppt":  _extract_pptx_metadata,
        ".txt":  _extract_txt_metadata,
        ".xlsx": _extract_excel_metadata,
        ".xls":  _extract_excel_metadata,
    }

    extractor = extractors.get(ext)
    if not extractor:
        logger.warning("No metadata extractor for extension: %s", ext)
        return {}

    try:
        metadata = extractor(filepath)
        # Always add file size
        metadata["file_size_kb"] = str(round(path.stat().st_size / 1024, 1))
        metadata["file_format"]  = ext.lstrip(".")
        _store_metadata(document_id, metadata)
        logger.info(
            "Stored %d metadata keys for document_id=%d.", len(metadata), document_id
        )
        return metadata
    except Exception as exc:
        logger.warning("Metadata extraction failed for %s: %s", filepath, exc)
        return {}


def _extract_pdf_metadata(filepath: str) -> dict:
    from langchain_community.document_loaders import PyPDFLoader
    loader = PyPDFLoader(filepath)
    pages  = loader.load()
    text   = " ".join(p.page_content for p in pages)
    return {
        "page_count": str(len(pages)),
        "word_count": str(len(text.split())),
    }


def _extract_docx_metadata(filepath: str) -> dict:
    from docx import Document
    doc        = Document(filepath)
    word_count = sum(len(p.text.split()) for p in doc.paragraphs)
    # Approximate page count: ~250 words per page
    page_count = max(1, round(word_count / 250))
    return {
        "page_count":  str(page_count),
        "word_count":  str(word_count),
        "table_count": str(len(doc.tables)),
    }


def _extract_pptx_metadata(filepath: str) -> dict:
    from pptx import Presentation
    prs      = Presentation(filepath)
    has_imgs = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
                has_imgs = True
                break
    return {
        "slide_count": str(len(prs.slides)),
        "has_images":  str(has_imgs),
    }


def _extract_txt_metadata(filepath: str) -> dict:
    with open(filepath, encoding="utf-8", errors="replace") as f:
        content = f.read()
    lines = content.splitlines()
    return {
        "line_count": str(len(lines)),
        "word_count": str(len(content.split())),
    }


def _extract_excel_metadata(filepath: str) -> dict:
    """
    Extract sheet names, row counts, and column names from an Excel file.
    Returns multiple metadata keys — one per sheet for rows and columns.
    """
    import openpyxl
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)

    metadata     = {}
    sheet_names  = wb.sheetnames
    metadata["sheet_names"] = ", ".join(sheet_names)
    metadata["sheet_count"] = str(len(sheet_names))

    for sheet_name in sheet_names:
        ws   = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # First row is assumed to be the header
        headers = [str(c) if c is not None else "" for c in rows[0]]
        data_rows = len(rows) - 1   # exclude header

        safe_name = sheet_name.replace(" ", "_")[:30]
        metadata[f"sheet_{safe_name}_columns"] = ", ".join(h for h in headers if h)
        metadata[f"sheet_{safe_name}_row_count"] = str(data_rows)

    wb.close()
    return metadata


# ---------------------------------------------------------------------------
# Storage and retrieval
# ---------------------------------------------------------------------------

def _store_metadata(document_id: int, metadata: dict) -> None:
    """Upsert all key-value metadata pairs for a document."""
    from services.database_service import _get_connection

    if not metadata:
        return

    with _get_connection() as conn:
        with conn.cursor() as cur:
            for key, value in metadata.items():
                cur.execute("""
                    INSERT INTO document_metadata (document_id, key, value)
                    VALUES (%s, %s, %s)
                    ON CONFLICT (document_id, key) DO UPDATE SET value = EXCLUDED.value;
                """, (document_id, key, str(value)))


def get_metadata(document_id: int) -> dict:
    """Return all metadata for a document as a dict."""
    import psycopg2.extras
    from services.database_service import _get_connection

    with _get_connection() as conn:
        with conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor) as cur:
            cur.execute(
                "SELECT key, value FROM document_metadata WHERE document_id = %s;",
                (document_id,)
            )
            rows = cur.fetchall()

    return {r["key"]: r["value"] for r in rows}


def get_all_metadata() -> list:
    """
    Return metadata for all documents joined with document info.
    Returns list of dicts: {document_id, filename, key, value}
    """
    import psycopg2.extras
    from services.database_service import _get_connection

    with _get_connection() as conn:
        with conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor) as cur:
            cur.execute("""
                SELECT d.id AS document_id, d.filename, m.key, m.value
                FROM document_metadata m
                JOIN documents d ON d.id = m.document_id
                ORDER BY d.id, m.key;
            """)
            return [dict(r) for r in cur.fetchall()]


def answer_metadata_question(query: str) -> dict:
    """
    Answer a metadata question by querying stored document metadata.

    Returns:
        dict with:
            answered : bool
            response : str    — human-readable answer
            data     : list   — raw metadata rows for UI display
    """
    from services.database_service import get_all_documents

    docs = get_all_documents()
    if not docs:
        return {
            "answered": True,
            "response": "No documents have been uploaded yet.",
            "data":     []
        }

    query_lower = query.lower()
    results     = []

    for doc in docs:
        metadata = get_metadata(doc["id"])
        if not metadata:
            continue

        # Filter metadata keys relevant to the query
        relevant = {}
        if any(w in query_lower for w in ["page", "slide"]):
            for k in ["page_count", "slide_count"]:
                if k in metadata:
                    relevant[k] = metadata[k]

        elif any(w in query_lower for w in ["row", "record"]):
            relevant = {k: v for k, v in metadata.items() if "row_count" in k}

        elif any(w in query_lower for w in ["column", "field", "header"]):
            relevant = {k: v for k, v in metadata.items() if "column" in k}

        elif any(w in query_lower for w in ["sheet"]):
            for k in ["sheet_names", "sheet_count"]:
                if k in metadata:
                    relevant[k] = metadata[k]

        elif any(w in query_lower for w in ["word"]):
            if "word_count" in metadata:
                relevant["word_count"] = metadata["word_count"]

        else:
            # General metadata question — return everything
            relevant = metadata

        if relevant:
            results.append({
                "filename": doc["filename"],
                "metadata": relevant,
            })

    if not results:
        return {
            "answered": False,
            "response": "",
            "data":     []
        }

    # Build human-readable response
    lines = []
    for r in results:
        lines.append(f"**{r['filename']}**")
        for k, v in r["metadata"].items():
            label = k.replace("_", " ").title()
            lines.append(f"  {label}: {v}")

    return {
        "answered": True,
        "response": "\n".join(lines),
        "data":     results,
    }

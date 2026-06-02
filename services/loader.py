"""
services/loader.py
------------------
Document text extraction with section heading markers.

Each loader embeds [HEADING: ...] markers into the extracted text wherever
a section heading is detected. The chunker uses these markers to prefix each
chunk with its section context — so the LLM always knows what section a chunk
belongs to, even after splitting.

Example output:
    [HEADING: Products]
    -CMS software -CRM software -ERP software...
    [HEADING: About Us]
    8Queens is an Enterprise Consulting...

Supported formats:
  PDF   (.pdf)  → PyPDFLoader  (text-based only)
  TXT   (.txt)  → TextLoader
  PPTX  (.pptx) → python-pptx  (slide titles as headings)
  PPT   (.ppt)  → python-pptx
  DOCX  (.docx) → python-docx  (Heading styles detected)
"""

import logging
import re
from pathlib import Path

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import SUPPORTED_EXTENSIONS

logger = logging.getLogger(__name__)

HEADING_MARKER = "[HEADING: {}]"


# ---------------------------------------------------------------------------
# Heading detection helpers
# ---------------------------------------------------------------------------

def _looks_like_heading(line: str) -> bool:
    """
    Heuristic heading detection for plain text and PDF output.
    A line is likely a heading if it is:
      - Short (under 80 characters)
      - Does not end with a period or comma
      - Is not purely numeric
      - Has at least 2 non-whitespace characters
    """
    line = line.strip()
    if not line or len(line) < 2 or len(line) > 80:
        return False
    if line.endswith(('.', ',', ';', ')', ']')):
        return False
    if re.match(r'^\d+\.?\d*$', line):
        return False
    # Title case or ALL CAPS lines are strong heading signals
    words = line.split()
    if len(words) == 0:
        return False
    title_case = sum(1 for w in words if w and w[0].isupper()) / len(words) >= 0.6
    all_caps   = line.upper() == line and any(c.isalpha() for c in line)
    ends_colon = line.endswith(':')
    return title_case or all_caps or ends_colon


def _inject_headings_from_lines(text: str) -> str:
    """
    Scan plain text line by line and inject [HEADING: ...] markers
    before lines that look like section headings.
    Used for PDF and TXT output.
    """
    lines  = text.split('\n')
    result = []

    for line in lines:
        stripped = line.strip()
        if stripped and _looks_like_heading(stripped):
            result.append(HEADING_MARKER.format(stripped))
        result.append(line)

    return '\n'.join(result)


# ---------------------------------------------------------------------------
# Private loaders
# ---------------------------------------------------------------------------

def _load_pdf(filepath: str) -> str:
    """
    Extract text from a text-based PDF using PyPDFLoader.
    Injects heading markers using heuristic detection.
    """
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(filepath)
    pages  = loader.load()

    if not pages:
        raise ValueError(f"No text could be extracted from PDF: {filepath}")

    parts = []
    for page in pages:
        content = page.page_content.strip()
        if content:
            parts.append(_inject_headings_from_lines(content))

    text = "\n\n".join(parts)

    if not text.strip():
        raise ValueError(
            f"PDF appears to contain no extractable text: {filepath}\n"
            "Scanned PDFs are not supported."
        )

    logger.info("PDF loaded: %d pages, %d characters.", len(pages), len(text))
    return text


def _load_txt(filepath: str) -> str:
    """
    Extract text from a plain text file.
    Injects heading markers using heuristic detection.
    """
    from langchain_community.document_loaders import TextLoader

    try:
        loader = TextLoader(filepath, encoding="utf-8")
        docs   = loader.load()
    except UnicodeDecodeError:
        loader = TextLoader(filepath, autodetect_encoding=True)
        docs   = loader.load()

    if not docs:
        raise ValueError(f"No text could be extracted from TXT file: {filepath}")

    raw  = "\n\n".join(doc.page_content for doc in docs if doc.page_content.strip())
    text = _inject_headings_from_lines(raw)

    if not text.strip():
        raise ValueError(f"TXT file is empty: {filepath}")

    logger.info("TXT loaded: %d characters.", len(text))
    return text


def _load_pptx(filepath: str) -> str:
    """
    Extract text from a PowerPoint file using python-pptx.
    Slide titles are used as [HEADING: ...] markers.
    Non-title shapes contribute body text under the slide heading.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    try:
        prs = Presentation(filepath)
    except Exception as exc:
        raise ValueError(f"Could not open PowerPoint file: {filepath}\n{exc}")

    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_lines  = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Detect title placeholder
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )

            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                if is_title and not title_text:
                    title_text = line
                else:
                    body_lines.append(line)

        parts = []
        if title_text:
            parts.append(HEADING_MARKER.format(title_text))
        elif body_lines:
            parts.append(HEADING_MARKER.format(f"Slide {slide_num}"))

        parts.extend(body_lines)

        if parts:
            slide_texts.append("\n".join(parts))

    if not slide_texts:
        raise ValueError(
            f"PowerPoint contains no extractable text: {filepath}\n"
            "Slides with only images are not supported."
        )

    text = "\n\n".join(slide_texts)
    logger.info("PPTX loaded: %d slides, %d characters.", len(slide_texts), len(text))
    return text


def _load_docx(filepath: str) -> str:
    """
    Extract text from a Word document using python-docx.
    Paragraphs with Heading styles are emitted as [HEADING: ...] markers.
    Table cells are extracted row by row.
    """
    from docx import Document

    try:
        doc = Document(filepath)
    except Exception as exc:
        raise ValueError(f"Could not open Word document: {filepath}\n{exc}")

    sections = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""

        # Detect heading styles: "Heading 1", "Heading 2", "Heading 3", etc.
        if "Heading" in style_name or "Title" in style_name:
            sections.append(HEADING_MARKER.format(text))
        else:
            sections.append(text)

    # Extract table content
    for table in doc.tables:
        for row in table.rows:
            row_text = "  |  ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                sections.append(row_text)

    if not sections:
        raise ValueError(
            f"Word document contains no extractable text: {filepath}"
        )

    text = "\n".join(sections)
    logger.info("DOCX loaded: %d paragraphs/rows, %d characters.", len(sections), len(text))
    return text


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

_LOADERS = {
    ".pdf":  _load_pdf,
    ".txt":  _load_txt,
    ".pptx": _load_pptx,
    ".ppt":  _load_pptx,
    ".docx": _load_docx,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def load_document(filepath: str) -> str:
    """
    Extract all text from a document and return it as a single string.
    Heading markers ([HEADING: ...]) are embedded for chunk enrichment.

    Args:
        filepath: Path to the document file.

    Returns:
        Extracted text with embedded heading markers.

    Raises:
        FileNotFoundError : File does not exist.
        ValueError        : Unsupported extension or no extractable text.
    """
    path = Path(filepath)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    if not path.is_file():
        raise ValueError(f"Path is not a file: {filepath}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: '{ext}'\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    loader_fn = _LOADERS[ext]
    logger.info("Loading document: %s (type: %s)", path.name, ext)

    text = loader_fn(str(path.resolve()))

    logger.info(
        "Document loaded: '%s' — %d characters extracted.",
        path.name, len(text),
    )
    return text


def get_supported_extensions() -> set:
    return set(SUPPORTED_EXTENSIONS)